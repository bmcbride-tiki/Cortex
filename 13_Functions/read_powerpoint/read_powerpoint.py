# =============================================================================
# read_powerpoint.py
# -----------------------------------------------------------------------------
# WHAT THIS FILE DOES
#   A Function: reads a local .pptx file's slide text and returns it as
#   text, one `--- Slide N ---` section per slide.
#
# WHAT IT INTERACTS WITH
#   - `python-pptx` (`Presentation`), for the actual slide/shape text
#     extraction.
#   - `test_read_powerpoint.py`, this file's paired test.
#   - `core_router.py`/`workflow_engine.py`, which dispatch this Function
#     the same way as any Task (generic `09_Functions` category), passing
#     its `file_path` as a JSON payload.
# =============================================================================

import sys
import json
from pathlib import Path
from typing import Dict, Any


class ReadPowerpoint:
    """Reads a local .pptx file's slide text and returns it as text, one
    '--- Slide N ---' section per slide. Standalone-capable; also
    drag-and-droppable onto the Workflow Builder canvas (category
    09_Functions), dispatched via CoreRouter same as any other Task."""

    def run(self, file_path: str) -> Dict[str, Any]:
        from pptx import Presentation

        src = Path(file_path)
        if not src.exists():
            return {"success": False, "response": f"File not found: {src}"}

        prs = Presentation(str(src))
        sections = []
        for i, slide in enumerate(prs.slides, start=1):
            lines = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    lines.append(shape.text_frame.text.strip())
            sections.append(f"--- Slide {i} ---\n" + "\n".join(lines))

        return {"success": True, "text": "\n\n".join(sections)}


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(json.dumps({"success": False, "response": 'Expected one JSON payload arg: {"file_path": "..."}'}))
        sys.exit(1)

    try:
        params = json.loads(sys.argv[1])
        result = ReadPowerpoint().run(file_path=params.get("file_path", ""))
    except Exception as e:
        result = {"success": False, "response": f"read_powerpoint error: {e}"}

    print(json.dumps(result))
    if not result.get("success"):
        sys.exit(1)
