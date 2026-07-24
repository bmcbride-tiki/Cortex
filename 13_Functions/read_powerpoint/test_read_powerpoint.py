# =============================================================================
# test_read_powerpoint.py
# -----------------------------------------------------------------------------
# WHAT THIS FILE DOES
#   Checks that read_powerpoint.py's `run()` extracts real slide title/body
#   text from a real generated `.pptx`, and fails cleanly on a missing file.
#
# WHAT IT INTERACTS WITH
#   - `read_powerpoint.py`, the file under test.
#   - `python-pptx`, used here to generate a throwaway source `.pptx`.
# =============================================================================

import sys
import tempfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
from read_powerpoint import ReadPowerpoint


def test_reads_slide_text():
    with tempfile.TemporaryDirectory() as tmp:
        from pptx import Presentation
        src = Path(tmp) / "in.pptx"
        prs = Presentation()
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Slide One Title"
        slide.placeholders[1].text_frame.text = "Body text"
        prs.save(str(src))

        result = ReadPowerpoint().run(str(src))
        assert result["success"] is True
        assert "Slide One Title" in result["text"]
        assert "Body text" in result["text"]
        assert "--- Slide 1 ---" in result["text"]


def test_missing_file_fails_cleanly():
    result = ReadPowerpoint().run("/no/such/file.pptx")
    assert result["success"] is False


if __name__ == "__main__":
    test_reads_slide_text()
    test_missing_file_fails_cleanly()
    print("All read_powerpoint self-checks passed.")
