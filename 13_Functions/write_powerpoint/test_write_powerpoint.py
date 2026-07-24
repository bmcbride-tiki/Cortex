# =============================================================================
# test_write_powerpoint.py
# -----------------------------------------------------------------------------
# WHAT THIS FILE DOES
#   Checks that write_powerpoint.py's `run()` creates a real `.pptx` with
#   one slide per double-newline-separated text block, titled from each
#   block's first line.
#
# WHAT IT INTERACTS WITH
#   - `write_powerpoint.py`, the file under test.
#   - `python-pptx`, used here to read back and verify the generated file.
# =============================================================================

import sys
import tempfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
from write_powerpoint import WritePowerpoint


def test_creates_one_slide_per_block():
    with tempfile.TemporaryDirectory() as tmp:
        result = WritePowerpoint().run(
            text="Slide One\nBody one\n\nSlide Two\nBody two",
            output_dir=tmp,
            filename="out",
        )
        assert result["success"] is True
        out_path = Path(result["file_path"])
        assert out_path.exists()
        assert out_path.suffix == ".pptx"

        from pptx import Presentation
        prs = Presentation(str(out_path))
        assert len(prs.slides) == 2
        assert prs.slides[0].shapes.title.text == "Slide One"
        assert prs.slides[1].shapes.title.text == "Slide Two"


if __name__ == "__main__":
    test_creates_one_slide_per_block()
    print("All write_powerpoint self-checks passed.")
