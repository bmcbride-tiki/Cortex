# =============================================================================
# write_powerpoint.py
# -----------------------------------------------------------------------------
# Copyright 2025 Brian McBride at Tiki-1 Studio
# WHAT THIS FILE DOES
#   A Function: creates a new .pptx file with one slide per double-newline
#   -separated block of text (first line of each block becomes the slide
#   title, the rest becomes the body).
#
# WHAT IT INTERACTS WITH
#   - `python-pptx` (`Presentation`), for the actual slide creation.
#   - `test_write_powerpoint.py`, this file's paired test.
#   - `core_router.py`/`workflow_engine.py`, which dispatch this Function
#     the same way as any Task (generic `09_Functions` category), passing
#     its `text`/`output_dir`/`filename` as a JSON payload.
# =============================================================================

import sys
import json
import time
from pathlib import Path
from typing import Dict, Any


class WritePowerpoint:
    """Creates a new .pptx file with one slide per double-newline-separated block of
    text (first line of each block becomes the slide title, the rest becomes the body).
    Standalone-capable; also drag-and-droppable onto the Workflow Builder canvas
    (category 09_Functions), dispatched via CoreRouter same as any other Task."""

    def run(self, text: str, output_dir: str, filename: str = "") -> Dict[str, Any]:
        from pptx import Presentation

        out_dir = Path(output_dir)
        out_dir.mkdir(parents=True, exist_ok=True)
        filename = (filename or f"export_{int(time.time())}").strip()
        if not filename.endswith(".pptx"):
            filename += ".pptx"
        out_path = out_dir / filename

        prs = Presentation()
        title_and_content_layout = prs.slide_layouts[1]

        blocks = [b.strip() for b in text.split("\n\n") if b.strip()] or [text]
        for block in blocks:
            lines = block.split("\n")
            title_text = lines[0]
            body_text = "\n".join(lines[1:]).strip()

            slide = prs.slides.add_slide(title_and_content_layout)
            slide.shapes.title.text = title_text
            if body_text and len(slide.placeholders) > 1:
                slide.placeholders[1].text_frame.text = body_text

        prs.save(str(out_path))
        return {"success": True, "file_path": str(out_path)}


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(json.dumps({"success": False, "response": 'Expected one JSON payload arg: {"text": "...", "output_dir": "...", "filename": "..."}'}))
        sys.exit(1)

    try:
        params = json.loads(sys.argv[1])
        result = WritePowerpoint().run(
            text=params.get("text", ""),
            output_dir=params.get("output_dir", ""),
            filename=params.get("filename", ""),
        )
    except Exception as e:
        result = {"success": False, "response": f"write_powerpoint error: {e}"}

    print(json.dumps(result))
    if not result.get("success"):
        sys.exit(1)
