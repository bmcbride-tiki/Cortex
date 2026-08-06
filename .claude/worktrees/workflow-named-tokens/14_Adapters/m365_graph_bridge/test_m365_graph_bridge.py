# Copyright 2025 Brian McBride at Tiki-1 Studio
import sys
import tempfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
import m365_graph_bridge as m365


def test_list_files_returns_shaped_listing():
    result = m365.list_files("/Reports")
    assert "files" in result
    assert result["files"][0]["name"]


def test_download_file_writes_real_placeholder():
    with tempfile.TemporaryDirectory() as tmp:
        result = m365.download_file("/Reports/Budget.xlsx", tmp)
        local_path = Path(result["local_path"])
        assert local_path.exists()
        assert local_path.name == "Budget.xlsx"


def test_upload_file_validates_real_source():
    with tempfile.TemporaryDirectory() as tmp:
        real_file = Path(tmp) / "report.docx"
        real_file.write_text("content")
        result = m365.upload_file(str(real_file), "/Reports/report.docx")
        assert result["item_id"].startswith("item_")

        try:
            m365.upload_file(str(Path(tmp) / "missing.docx"), "/x")
            assert False, "expected FileNotFoundError for a missing local file"
        except FileNotFoundError:
            pass


def test_excel_range_get_and_set():
    get_result = m365.get_excel_range("Budget.xlsx", "Sheet1", "A1:B2")
    assert len(get_result["values"]) == 2

    set_result = m365.set_excel_range("Budget.xlsx", "Sheet1", "A1:B2", [["x", "y"]])
    assert set_result["row_count"] == 1


def test_powerbi_refresh_and_list():
    refresh_result = m365.refresh_powerbi_dataset("dataset_123")
    assert refresh_result["refresh_request_id"].startswith("refresh_")

    list_result = m365.list_powerbi_reports()
    assert len(list_result["reports"]) >= 1


def test_send_mail_splits_recipients():
    result = m365.send_mail("a@example.com, b@example.com", "Subject", "Body")
    assert result["to"] == ["a@example.com", "b@example.com"]


def test_list_messages_respects_top():
    result = m365.list_messages("inbox", top=1)
    assert len(result["messages"]) == 1


def test_calendar_event_round_trip():
    events = m365.list_calendar_events()
    assert len(events["events"]) >= 1

    created = m365.create_calendar_event("Sync", "2026-08-01T10:00:00Z", "2026-08-01T11:00:00Z", "a@example.com,b@example.com")
    assert created["event_id"].startswith("evt_")
    assert created["attendees"] == ["a@example.com", "b@example.com"]

    try:
        m365.create_calendar_event("", "", "", "")
        assert False, "expected ValueError for missing required fields"
    except ValueError:
        pass


def test_teams_and_channels():
    teams = m365.list_teams()
    assert len(teams["teams"]) >= 1

    channels = m365.list_channels(teams["teams"][0]["id"])
    assert len(channels["channels"]) >= 1

    try:
        m365.list_channels("")
        assert False, "expected ValueError for missing team_id"
    except ValueError:
        pass


def test_channel_and_chat_messaging():
    posted = m365.post_channel_message("team_1", "channel_1", "hello")
    assert posted["status"] == "posted (mock)"

    chat_list = m365.list_chat_messages("chat_1")
    assert len(chat_list["messages"]) >= 1

    sent = m365.send_chat_message("chat_1", "hi")
    assert sent["status"] == "sent (mock)"


def test_search_messages_filters_by_sender():
    result = m365.search_messages(query="curriculum", sender="registrar")
    assert len(result["messages"]) == 1
    assert "registrar" in result["messages"][0]["from"]

    try:
        m365.search_messages()
        assert False, "expected ValueError when neither query nor sender given"
    except ValueError:
        pass


def test_generate_pptx_from_word():
    with tempfile.TemporaryDirectory() as tmp:
        from docx import Document
        src = Path(tmp) / "source.docx"
        doc = Document()
        doc.add_paragraph("Section 1: Introduction")
        doc.add_paragraph("Body text.")
        doc.save(str(src))

        result = m365.generate_pptx_from_word(str(src), tmp)
        assert result["slide_count"] == 1
        out_path = Path(result["file_path"])
        assert out_path.exists()

        from pptx import Presentation
        prs = Presentation(str(out_path))
        assert "[MOCK Copilot-generated]" in prs.slides[0].shapes.title.text

    try:
        m365.generate_pptx_from_word("/no/such/file.docx", "/tmp")
        assert False, "expected FileNotFoundError for a missing word file"
    except FileNotFoundError:
        pass


def test_sharepoint_sites():
    all_sites = m365.list_sharepoint_sites()
    assert len(all_sites["sites"]) >= 2

    filtered = m365.list_sharepoint_sites(query="Exam")
    assert len(filtered["sites"]) == 1

    site = m365.get_sharepoint_site("example.sharepoint.com:/sites/apprenticeship")
    assert site["id"].startswith("site_")

    try:
        m365.get_sharepoint_site("")
        assert False, "expected ValueError for missing site_path"
    except ValueError:
        pass


def test_sharepoint_lists():
    lists = m365.list_sharepoint_lists("site_1")
    assert len(lists["lists"]) >= 1

    items = m365.list_sharepoint_list_items("site_1", "list_1")
    assert len(items["items"]) >= 1

    created = m365.create_sharepoint_list_item("site_1", "list_1", {"Title": "New Item"})
    assert created["item_id"].startswith("item_")

    try:
        m365.create_sharepoint_list_item("site_1", "list_1", {})
        assert False, "expected ValueError for empty fields"
    except ValueError:
        pass


def test_onenote():
    notebooks = m365.list_onenote_notebooks()
    assert len(notebooks["notebooks"]) >= 1

    pages = m365.list_onenote_pages(notebooks["notebooks"][0]["id"])
    assert len(pages["pages"]) >= 1

    content = m365.get_onenote_page_content(pages["pages"][0]["id"])
    assert "<html>" in content["content_html"]

    created = m365.create_onenote_page("section_1", "New Page", "<p>Body</p>")
    assert created["page_id"].startswith("page_")


def test_onedrive_extras():
    recent = m365.list_recent_onedrive_files()
    assert len(recent["files"]) >= 1

    link = m365.create_onedrive_sharing_link("/Reports/Budget.xlsx")
    assert link["share_url"].startswith("https://")

    try:
        m365.create_onedrive_sharing_link("")
        assert False, "expected ValueError for missing file_path"
    except ValueError:
        pass


def test_mock_mode_off_fails_clearly():
    m365.MOCK_MODE = False
    try:
        try:
            m365.list_files("/x")
            assert False, "expected RuntimeError when MOCK_MODE is off"
        except RuntimeError as e:
            assert "not configured" in str(e)
    finally:
        m365.MOCK_MODE = True


if __name__ == "__main__":
    test_list_files_returns_shaped_listing()
    test_download_file_writes_real_placeholder()
    test_upload_file_validates_real_source()
    test_excel_range_get_and_set()
    test_powerbi_refresh_and_list()
    test_send_mail_splits_recipients()
    test_list_messages_respects_top()
    test_calendar_event_round_trip()
    test_teams_and_channels()
    test_channel_and_chat_messaging()
    test_search_messages_filters_by_sender()
    test_generate_pptx_from_word()
    test_sharepoint_sites()
    test_sharepoint_lists()
    test_onenote()
    test_onedrive_extras()
    test_mock_mode_off_fails_clearly()
    print("All m365_graph_bridge self-checks passed.")
