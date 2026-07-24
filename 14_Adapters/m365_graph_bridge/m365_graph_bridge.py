# =============================================================================
# m365_graph_bridge.py
# -----------------------------------------------------------------------------
# WHAT THIS FILE DOES
#   Adapter for Microsoft 365 via the Microsoft Graph API (OneDrive/
#   SharePoint file access, Excel's Workbook API for cell-level operations)
#   plus the Power BI REST API (dataset refresh, report listing) -- both use
#   the same Azure AD/Entra ID app registration and MSAL auth flow, so they
#   live in one adapter rather than two. No Azure AD app registration exists
#   for this project yet, so every action runs in MOCK_MODE by default.
#
# WHAT IT INTERACTS WITH
#   - `core_router.py`, which is what actually runs this file when
#     `workflow_engine.py`'s M365/Power BI function nodes dispatch to it.
#   - Downstream: files this bridge downloads are meant to be handed to
#     existing content tools (`13_Functions/import_from_word`,
#     `read_powerpoint`, openpyxl-based tools, etc.) rather than parsed
#     here -- this adapter's job is only getting files in and out of
#     OneDrive/SharePoint, not reading their content.
#
# KEY FUNCTIONALITY NOTES
#   - Same MOCK_MODE / JSON-line contract as notebooklm_bridge.py.
#   - Power Query itself has no standalone public automation API -- the
#     real way to "run" a Power Query transformation programmatically is by
#     triggering a Power BI dataset refresh (refresh_powerbi_dataset below)
#     or an Excel workbook refresh, not a separate integration surface.
#   - download_file/upload_file still touch the real filesystem even in
#     mock mode (write a placeholder file / validate the source file
#     exists) so a workflow chaining into a real parsing step has something
#     real to work with, and a genuine usage mistake (bad local path) still
#     surfaces immediately -- same "realistic data contracts" principle as
#     notebooklm_bridge.py.
# =============================================================================

# 14_Adapters/m365_graph_bridge/m365_graph_bridge.py
import sys
sys.dont_write_bytecode = True

import os
import json
import uuid
import argparse
from pathlib import Path
from typing import Any, Dict, List

MOCK_MODE = os.environ.get("M365_MOCK_MODE", "1") != "0"

NOT_CONFIGURED_MESSAGE = (
    "M365/Graph API access is not configured yet. Set M365_MOCK_MODE=1 (the default) "
    "to use simulated responses, or register an Azure AD app (tenant ID, client ID, "
    "client secret/certificate, Graph + Power BI API permissions) and wire up MSAL + "
    "the msgraph-sdk inside m365_graph_bridge.py."
)


def list_files(folder_path: str) -> Dict[str, Any]:
    if not MOCK_MODE:
        raise RuntimeError(NOT_CONFIGURED_MESSAGE)
    folder_path = folder_path or "/"
    files = [
        {"name": "Quarterly Report.docx", "item_id": f"item_{uuid.uuid4().hex[:10]}", "size": 48213, "web_url": f"https://example.sharepoint.com{folder_path}/Quarterly%20Report.docx"},
        {"name": "Budget.xlsx", "item_id": f"item_{uuid.uuid4().hex[:10]}", "size": 102934, "web_url": f"https://example.sharepoint.com{folder_path}/Budget.xlsx"},
    ]
    return {"files": files}


def download_file(file_path: str, local_output_dir: str) -> Dict[str, Any]:
    if not MOCK_MODE:
        raise RuntimeError(NOT_CONFIGURED_MESSAGE)
    if not file_path:
        raise ValueError("download_file requires a file_path.")

    out_dir = Path(local_output_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    local_path = out_dir / Path(file_path).name
    local_path.write_text(f"[MOCK M365 download] placeholder content for {file_path}", encoding="utf-8")
    return {"local_path": str(local_path)}


def upload_file(local_path: str, destination_path: str) -> Dict[str, Any]:
    if not MOCK_MODE:
        raise RuntimeError(NOT_CONFIGURED_MESSAGE)
    if not local_path:
        raise ValueError("upload_file requires a local_path.")

    src = Path(local_path)
    if not src.exists():
        raise FileNotFoundError(f"Local file not found: {src}")

    return {"item_id": f"item_{uuid.uuid4().hex[:10]}", "web_url": f"https://example.sharepoint.com{destination_path or '/' + src.name}"}


def get_excel_range(file_path: str, worksheet: str, range_address: str) -> Dict[str, Any]:
    if not MOCK_MODE:
        raise RuntimeError(NOT_CONFIGURED_MESSAGE)
    if not file_path or not worksheet or not range_address:
        raise ValueError("get_excel_range requires file_path, worksheet, and range_address.")
    # Realistic-shaped placeholder: Graph's Workbook API returns a 2D array of cell values.
    return {"values": [["Header A", "Header B"], ["Row1 A", "Row1 B"]]}


def set_excel_range(file_path: str, worksheet: str, range_address: str, values: List[List[Any]]) -> Dict[str, Any]:
    if not MOCK_MODE:
        raise RuntimeError(NOT_CONFIGURED_MESSAGE)
    if not file_path or not worksheet or not range_address:
        raise ValueError("set_excel_range requires file_path, worksheet, and range_address.")
    return {"updated_range": range_address, "row_count": len(values or [])}


def refresh_powerbi_dataset(dataset_id: str) -> Dict[str, Any]:
    if not MOCK_MODE:
        raise RuntimeError(NOT_CONFIGURED_MESSAGE)
    if not dataset_id:
        raise ValueError("refresh_powerbi_dataset requires a dataset_id.")
    return {"refresh_request_id": f"refresh_{uuid.uuid4().hex[:10]}", "status": "Unknown (mock -- would be 'InProgress' immediately after a real trigger)"}


def list_powerbi_reports(workspace_id: str = "") -> Dict[str, Any]:
    if not MOCK_MODE:
        raise RuntimeError(NOT_CONFIGURED_MESSAGE)
    reports = [
        {"id": f"report_{uuid.uuid4().hex[:10]}", "name": "Regional Enrollment Trends", "dataset_id": f"dataset_{uuid.uuid4().hex[:10]}"},
    ]
    return {"reports": reports}


# --- Outlook ---------------------------------------------------------------------

def send_mail(to: str, subject: str, body: str) -> Dict[str, Any]:
    if not MOCK_MODE:
        raise RuntimeError(NOT_CONFIGURED_MESSAGE)
    if not to:
        raise ValueError("send_mail requires a to address (comma-separated for multiple).")
    return {"message_id": f"msg_{uuid.uuid4().hex[:10]}", "to": [a.strip() for a in to.split(",") if a.strip()], "status": "sent (mock)"}


def list_messages(folder: str = "inbox", top: int = 10) -> Dict[str, Any]:
    if not MOCK_MODE:
        raise RuntimeError(NOT_CONFIGURED_MESSAGE)
    messages = [
        {"id": f"msg_{uuid.uuid4().hex[:10]}", "from": "sme.lead@example.com", "subject": "Curriculum review sign-off", "received": "2026-07-20T09:15:00Z", "preview": "Please review the attached section before Friday..."},
        {"id": f"msg_{uuid.uuid4().hex[:10]}", "from": "registrar@example.com", "subject": "Exam schedule confirmed", "received": "2026-07-19T14:02:00Z", "preview": "The Q3 exam schedule has been finalized..."},
    ]
    return {"folder": folder or "inbox", "messages": messages[: max(top, 0) or len(messages)]}


def search_messages(query: str = "", sender: str = "", folder: str = "inbox", top: int = 10) -> Dict[str, Any]:
    # Mirrors Graph's real /me/messages?$search=/$filter=from/emailAddress/address
    # capability -- a keyword/topic search plus an optional sender filter, unlike
    # list_messages' plain unfiltered folder listing.
    if not MOCK_MODE:
        raise RuntimeError(NOT_CONFIGURED_MESSAGE)
    if not query and not sender:
        raise ValueError("search_messages requires a query and/or a sender to filter by.")
    candidate_messages = [
        {"id": f"msg_{uuid.uuid4().hex[:10]}", "from": "sme.lead@example.com", "subject": "Curriculum review sign-off", "received": "2026-07-20T09:15:00Z", "preview": f"[MOCK match for query={query!r}] Please review the attached section before Friday..."},
        {"id": f"msg_{uuid.uuid4().hex[:10]}", "from": "registrar@example.com", "subject": "Exam schedule confirmed", "received": "2026-07-19T14:02:00Z", "preview": f"[MOCK match for query={query!r}] The Q3 exam schedule has been finalized..."},
    ]
    if sender:
        candidate_messages = [m for m in candidate_messages if sender.lower() in m["from"].lower()]
    return {"folder": folder or "inbox", "query": query, "sender": sender, "messages": candidate_messages[: max(top, 0) or len(candidate_messages)]}


def list_calendar_events(start_date: str = "", end_date: str = "") -> Dict[str, Any]:
    if not MOCK_MODE:
        raise RuntimeError(NOT_CONFIGURED_MESSAGE)
    events = [
        {"id": f"evt_{uuid.uuid4().hex[:10]}", "subject": "Curriculum Committee Sync", "start": start_date or "2026-07-23T10:00:00Z", "end": end_date or "2026-07-23T11:00:00Z"},
    ]
    return {"events": events}


def create_calendar_event(subject: str, start: str, end: str, attendees: str = "") -> Dict[str, Any]:
    if not MOCK_MODE:
        raise RuntimeError(NOT_CONFIGURED_MESSAGE)
    if not subject or not start or not end:
        raise ValueError("create_calendar_event requires subject, start, and end.")
    return {
        "event_id": f"evt_{uuid.uuid4().hex[:10]}",
        "subject": subject,
        "start": start,
        "end": end,
        "attendees": [a.strip() for a in attendees.split(",") if a.strip()],
    }


# --- Teams -------------------------------------------------------------------------

def list_teams() -> Dict[str, Any]:
    if not MOCK_MODE:
        raise RuntimeError(NOT_CONFIGURED_MESSAGE)
    teams = [
        {"id": f"team_{uuid.uuid4().hex[:10]}", "name": "Apprenticeship Curriculum"},
        {"id": f"team_{uuid.uuid4().hex[:10]}", "name": "Exam Standards Board"},
    ]
    return {"teams": teams}


def list_channels(team_id: str) -> Dict[str, Any]:
    if not MOCK_MODE:
        raise RuntimeError(NOT_CONFIGURED_MESSAGE)
    if not team_id:
        raise ValueError("list_channels requires a team_id.")
    channels = [
        {"id": f"channel_{uuid.uuid4().hex[:10]}", "name": "General"},
        {"id": f"channel_{uuid.uuid4().hex[:10]}", "name": "Curriculum Reviews"},
    ]
    return {"team_id": team_id, "channels": channels}


def post_channel_message(team_id: str, channel_id: str, message: str) -> Dict[str, Any]:
    if not MOCK_MODE:
        raise RuntimeError(NOT_CONFIGURED_MESSAGE)
    if not team_id or not channel_id or not message:
        raise ValueError("post_channel_message requires team_id, channel_id, and message.")
    return {"message_id": f"chanmsg_{uuid.uuid4().hex[:10]}", "status": "posted (mock)"}


def list_chat_messages(chat_id: str) -> Dict[str, Any]:
    if not MOCK_MODE:
        raise RuntimeError(NOT_CONFIGURED_MESSAGE)
    if not chat_id:
        raise ValueError("list_chat_messages requires a chat_id.")
    messages = [
        {"id": f"chatmsg_{uuid.uuid4().hex[:10]}", "from": "colleague@example.com", "sent": "2026-07-21T16:40:00Z", "text": "Draft looks good, one edit on section 3."},
    ]
    return {"chat_id": chat_id, "messages": messages}


def send_chat_message(chat_id: str, message: str) -> Dict[str, Any]:
    if not MOCK_MODE:
        raise RuntimeError(NOT_CONFIGURED_MESSAGE)
    if not chat_id or not message:
        raise ValueError("send_chat_message requires chat_id and message.")
    return {"message_id": f"chatmsg_{uuid.uuid4().hex[:10]}", "status": "sent (mock)"}


# --- M365 Copilot (classified as a Copilot connection, not plain M365, even
# though the code lives here alongside the file it operates on -- see the
# Task wrapper's explicit model="copilot" tag in server.py) ---

def generate_pptx_from_word(word_file_path: str, output_dir: str, filename: str = "") -> Dict[str, Any]:
    # Honest limitation: this mocks the FILE-HANDLING side of M365 Copilot's real
    # "Create presentation from file" feature (no such API exists to call, mock or
    # otherwise) with a straightforward mechanical Word -> PowerPoint conversion --
    # every 5 non-empty paragraphs becomes one slide, first paragraph as its title.
    # It does NOT use AI to restructure/summarize content the way real Copilot would.
    # Swap in a real Copilot API call here once one exists; everything downstream
    # (the Task wrapper, the workflow node) stays the same.
    if not MOCK_MODE:
        raise RuntimeError(NOT_CONFIGURED_MESSAGE)
    if not word_file_path:
        raise ValueError("generate_pptx_from_word requires a word_file_path.")

    from docx import Document
    from pptx import Presentation

    src = Path(word_file_path)
    if not src.exists():
        raise FileNotFoundError(f"Word file not found: {src}")

    doc = Document(str(src))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    if not paragraphs:
        raise ValueError(f"No text found in {src} to build slides from.")

    out_dir = Path(output_dir) if output_dir else src.parent
    out_dir.mkdir(parents=True, exist_ok=True)
    filename = (filename or f"{src.stem}_generated").strip()
    if not filename.endswith(".pptx"):
        filename += ".pptx"
    out_path = out_dir / filename

    prs = Presentation()
    layout = prs.slide_layouts[1]
    chunk_size = 5
    slide_count = 0
    for i in range(0, len(paragraphs), chunk_size):
        chunk = paragraphs[i:i + chunk_size]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = f"[MOCK Copilot-generated] {chunk[0]}"
        body = "\n".join(chunk[1:])
        if body and len(slide.placeholders) > 1:
            slide.placeholders[1].text_frame.text = body
        slide_count += 1
    prs.save(str(out_path))

    return {"file_path": str(out_path), "slide_count": slide_count}


# --- SharePoint Online sites ---------------------------------------------------------

def list_sharepoint_sites(query: str = "") -> Dict[str, Any]:
    if not MOCK_MODE:
        raise RuntimeError(NOT_CONFIGURED_MESSAGE)
    sites = [
        {"id": f"site_{uuid.uuid4().hex[:10]}", "name": "Apprenticeship Program", "web_url": "https://example.sharepoint.com/sites/apprenticeship"},
        {"id": f"site_{uuid.uuid4().hex[:10]}", "name": "Exam Standards", "web_url": "https://example.sharepoint.com/sites/examstandards"},
    ]
    if query:
        sites = [s for s in sites if query.lower() in s["name"].lower()]
    return {"query": query, "sites": sites}


def get_sharepoint_site(site_path: str) -> Dict[str, Any]:
    if not MOCK_MODE:
        raise RuntimeError(NOT_CONFIGURED_MESSAGE)
    if not site_path:
        raise ValueError("get_sharepoint_site requires a site_path (e.g. 'example.sharepoint.com:/sites/apprenticeship').")
    return {
        "id": f"site_{uuid.uuid4().hex[:10]}",
        "name": Path(site_path.rstrip("/")).name or site_path,
        "web_url": f"https://{site_path.lstrip('/')}" if not site_path.startswith("http") else site_path,
    }


# --- Microsoft Lists -----------------------------------------------------------------

def list_sharepoint_lists(site_id: str) -> Dict[str, Any]:
    if not MOCK_MODE:
        raise RuntimeError(NOT_CONFIGURED_MESSAGE)
    if not site_id:
        raise ValueError("list_sharepoint_lists requires a site_id.")
    lists = [
        {"id": f"list_{uuid.uuid4().hex[:10]}", "name": "Apprentice Intake Tracker"},
        {"id": f"list_{uuid.uuid4().hex[:10]}", "name": "Curriculum Review Status"},
    ]
    return {"site_id": site_id, "lists": lists}


def list_sharepoint_list_items(site_id: str, list_id: str) -> Dict[str, Any]:
    if not MOCK_MODE:
        raise RuntimeError(NOT_CONFIGURED_MESSAGE)
    if not site_id or not list_id:
        raise ValueError("list_sharepoint_list_items requires site_id and list_id.")
    items = [
        {"id": f"item_{uuid.uuid4().hex[:10]}", "fields": {"Title": "Electrician - Period 3", "Status": "In Review"}},
        {"id": f"item_{uuid.uuid4().hex[:10]}", "fields": {"Title": "Plumber - Period 2", "Status": "Approved"}},
    ]
    return {"site_id": site_id, "list_id": list_id, "items": items}


def create_sharepoint_list_item(site_id: str, list_id: str, fields: Dict[str, Any]) -> Dict[str, Any]:
    if not MOCK_MODE:
        raise RuntimeError(NOT_CONFIGURED_MESSAGE)
    if not site_id or not list_id:
        raise ValueError("create_sharepoint_list_item requires site_id and list_id.")
    if not isinstance(fields, dict) or not fields:
        raise ValueError("create_sharepoint_list_item requires a non-empty fields object.")
    return {"item_id": f"item_{uuid.uuid4().hex[:10]}", "fields": fields}


# --- OneNote ---------------------------------------------------------------------

def list_onenote_notebooks() -> Dict[str, Any]:
    if not MOCK_MODE:
        raise RuntimeError(NOT_CONFIGURED_MESSAGE)
    notebooks = [
        {"id": f"notebook_{uuid.uuid4().hex[:10]}", "name": "Curriculum Committee Notes"},
        {"id": f"notebook_{uuid.uuid4().hex[:10]}", "name": "Exam Standards Board"},
    ]
    return {"notebooks": notebooks}


def list_onenote_pages(notebook_id: str) -> Dict[str, Any]:
    if not MOCK_MODE:
        raise RuntimeError(NOT_CONFIGURED_MESSAGE)
    if not notebook_id:
        raise ValueError("list_onenote_pages requires a notebook_id.")
    pages = [
        {"id": f"page_{uuid.uuid4().hex[:10]}", "title": "2026-07-20 Meeting Notes"},
        {"id": f"page_{uuid.uuid4().hex[:10]}", "title": "Action Items"},
    ]
    return {"notebook_id": notebook_id, "pages": pages}


def get_onenote_page_content(page_id: str) -> Dict[str, Any]:
    if not MOCK_MODE:
        raise RuntimeError(NOT_CONFIGURED_MESSAGE)
    if not page_id:
        raise ValueError("get_onenote_page_content requires a page_id.")
    # Real Graph OneNote pages return HTML content -- mirrored here for realism.
    return {"page_id": page_id, "content_html": f"<html><body><p>[MOCK OneNote content for page {page_id}]</p></body></html>"}


def create_onenote_page(section_id: str, title: str, content: str) -> Dict[str, Any]:
    if not MOCK_MODE:
        raise RuntimeError(NOT_CONFIGURED_MESSAGE)
    if not section_id or not title:
        raise ValueError("create_onenote_page requires section_id and title.")
    return {"page_id": f"page_{uuid.uuid4().hex[:10]}", "title": title, "status": "created (mock)"}


# --- OneDrive (beyond list_files/download_file/upload_file above) --------------------

def list_recent_onedrive_files() -> Dict[str, Any]:
    if not MOCK_MODE:
        raise RuntimeError(NOT_CONFIGURED_MESSAGE)
    files = [
        {"name": "Quarterly Report.docx", "item_id": f"item_{uuid.uuid4().hex[:10]}", "last_modified": "2026-07-21T10:05:00Z"},
        {"name": "Budget.xlsx", "item_id": f"item_{uuid.uuid4().hex[:10]}", "last_modified": "2026-07-20T15:40:00Z"},
    ]
    return {"files": files}


def create_onedrive_sharing_link(file_path: str) -> Dict[str, Any]:
    if not MOCK_MODE:
        raise RuntimeError(NOT_CONFIGURED_MESSAGE)
    if not file_path:
        raise ValueError("create_onedrive_sharing_link requires a file_path.")
    return {"share_url": f"https://example.sharepoint.com/:x:/g/personal/mock/{uuid.uuid4().hex[:16]}", "link_type": "view"}


def main():
    # Expects one command-line argument: a small JSON instruction, e.g.
    #   {"action": "list_files", "folder_path": "/Reports"}
    # Always prints back exactly one line of JSON: either
    # {"success": true, ...} or {"success": false, "response": "<error>"}.
    parser = argparse.ArgumentParser(description="Workbrain M365/Graph + Power BI Integration Bridge (mock-mode until an Azure AD app registration exists)")
    parser.add_argument("payload", nargs="?", default="", help="Inline JSON parameter from router")
    args = parser.parse_args()

    if not args.payload:
        print(json.dumps({
            "success": False,
            "response": 'No JSON payload provided. Expected: {"action": "list_files|download_file|upload_file|get_excel_range|set_excel_range|refresh_powerbi_dataset|list_powerbi_reports", ...}',
        }))
        sys.exit(1)

    try:
        params = json.loads(args.payload)
    except json.JSONDecodeError as e:
        print(json.dumps({"success": False, "response": f"Malformed JSON payload: {e}"}))
        sys.exit(1)

    action = params.get("action", "")

    try:
        if action == "list_files":
            result = {"success": True, **list_files(params.get("folder_path", ""))}
        elif action == "download_file":
            result = {"success": True, **download_file(params.get("file_path", ""), params.get("local_output_dir", ""))}
        elif action == "upload_file":
            result = {"success": True, **upload_file(params.get("local_path", ""), params.get("destination_path", ""))}
        elif action == "get_excel_range":
            result = {"success": True, **get_excel_range(params.get("file_path", ""), params.get("worksheet", ""), params.get("range_address", ""))}
        elif action == "set_excel_range":
            result = {"success": True, **set_excel_range(params.get("file_path", ""), params.get("worksheet", ""), params.get("range_address", ""), params.get("values", []))}
        elif action == "refresh_powerbi_dataset":
            result = {"success": True, **refresh_powerbi_dataset(params.get("dataset_id", ""))}
        elif action == "list_powerbi_reports":
            result = {"success": True, **list_powerbi_reports(params.get("workspace_id", ""))}
        elif action == "send_mail":
            result = {"success": True, **send_mail(params.get("to", ""), params.get("subject", ""), params.get("body", ""))}
        elif action == "list_messages":
            result = {"success": True, **list_messages(params.get("folder", "inbox"), int(params.get("top") or 10))}
        elif action == "search_messages":
            result = {"success": True, **search_messages(params.get("query", ""), params.get("sender", ""), params.get("folder", "inbox"), int(params.get("top") or 10))}
        elif action == "list_calendar_events":
            result = {"success": True, **list_calendar_events(params.get("start_date", ""), params.get("end_date", ""))}
        elif action == "create_calendar_event":
            result = {"success": True, **create_calendar_event(params.get("subject", ""), params.get("start", ""), params.get("end", ""), params.get("attendees", ""))}
        elif action == "list_teams":
            result = {"success": True, **list_teams()}
        elif action == "list_channels":
            result = {"success": True, **list_channels(params.get("team_id", ""))}
        elif action == "post_channel_message":
            result = {"success": True, **post_channel_message(params.get("team_id", ""), params.get("channel_id", ""), params.get("message", ""))}
        elif action == "list_chat_messages":
            result = {"success": True, **list_chat_messages(params.get("chat_id", ""))}
        elif action == "send_chat_message":
            result = {"success": True, **send_chat_message(params.get("chat_id", ""), params.get("message", ""))}
        elif action == "generate_pptx_from_word":
            result = {"success": True, **generate_pptx_from_word(params.get("word_file_path", ""), params.get("output_dir", ""), params.get("filename", ""))}
        elif action == "list_sharepoint_sites":
            result = {"success": True, **list_sharepoint_sites(params.get("query", ""))}
        elif action == "get_sharepoint_site":
            result = {"success": True, **get_sharepoint_site(params.get("site_path", ""))}
        elif action == "list_sharepoint_lists":
            result = {"success": True, **list_sharepoint_lists(params.get("site_id", ""))}
        elif action == "list_sharepoint_list_items":
            result = {"success": True, **list_sharepoint_list_items(params.get("site_id", ""), params.get("list_id", ""))}
        elif action == "create_sharepoint_list_item":
            result = {"success": True, **create_sharepoint_list_item(params.get("site_id", ""), params.get("list_id", ""), params.get("fields", {}))}
        elif action == "list_onenote_notebooks":
            result = {"success": True, **list_onenote_notebooks()}
        elif action == "list_onenote_pages":
            result = {"success": True, **list_onenote_pages(params.get("notebook_id", ""))}
        elif action == "get_onenote_page_content":
            result = {"success": True, **get_onenote_page_content(params.get("page_id", ""))}
        elif action == "create_onenote_page":
            result = {"success": True, **create_onenote_page(params.get("section_id", ""), params.get("title", ""), params.get("content", ""))}
        elif action == "list_recent_onedrive_files":
            result = {"success": True, **list_recent_onedrive_files()}
        elif action == "create_onedrive_sharing_link":
            result = {"success": True, **create_onedrive_sharing_link(params.get("file_path", ""))}
        else:
            result = {"success": False, "response": f"Unknown action: {action}"}
    except Exception as e:
        result = {"success": False, "response": f"M365 Graph bridge error: {e}"}

    print(json.dumps(result))
    if not result.get("success"):
        sys.exit(1)


if __name__ == "__main__":
    main()
